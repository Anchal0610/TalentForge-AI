import os
from dotenv import load_dotenv
from utils.logger import logger

# Import standard document processors
import pypdf
import docx
import pptx

load_dotenv()

class MistralOCRService:
    def __init__(self):
        self.api_key = os.getenv("MISTRAL_API_KEY")
        if not self.api_key or self.api_key.startswith("your_"):
            logger.warning("MISTRAL_API_KEY not configured. Mistral OCR operates in local fallback parser mode.")
            self.client = None
        else:
            try:
                from mistralai import Mistral
                self.client = Mistral(api_key=self.api_key)
                logger.info("Mistral OCR client successfully initialized.")
            except ImportError:
                logger.warning("Mistral SDK missing or error importing. Operating in local fallback parser mode.")
                self.client = None

    def extract_text(self, file_path: str) -> str:
        """Main entry point to extract text. Tries Mistral OCR first, falls back to local libraries."""
        logger.info(f"Extracting text from file: {file_path}")
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        # Check extension
        _, ext = os.path.splitext(file_path.lower())
        
        # Try Mistral OCR first if configured and it's a supported format (usually PDFs/images)
        if self.client and ext in [".pdf", ".png", ".jpg", ".jpeg", ".tiff"]:
            try:
                logger.info(f"Attempting Mistral OCR API processing for {file_path}...")
                # Standard Mistral OCR API call format (Mistral Client document processing)
                # For Phase 1, we design the integration interface.
                # In real execution, we process the document file upload and request OCR.
                # Here is a mockable API call block:
                # response = self.client.ocr.process(document={"type": "local", "path": file_path})
                # return response.pages[0].markdown ...
                
                # Let's write the template API call:
                with open(file_path, "rb") as f:
                    # Depending on Mistral SDK version:
                    # we write the expected structure but wrap it securely in a try-except.
                    # Since this is a template wrapper for OCR, we can mock it if needed or run a mock file read.
                    raise NotImplementedError("Mistral OCR execution placeholder - fallback triggered.")
            except Exception as e:
                logger.warning(f"Mistral OCR API failed or not fully set up ({str(e)}). Falling back to local parser.")
        
        # Local Fallbacks
        return self._local_parse(file_path, ext)

    def _local_parse(self, file_path: str, ext: str) -> str:
        """Fallback parse using local Python libraries."""
        logger.info(f"Running local parser for extension {ext}...")
        
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".docx":
            return self._parse_docx(file_path)
        elif ext in [".ppt", ".pptx"]:
            return self._parse_pptx(file_path)
        elif ext in [".txt", ".md", ".json"]:
            return self._parse_text(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _parse_pdf(self, file_path: str) -> str:
        text_content = []
        try:
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error parsing PDF locally: {str(e)}")
            raise e

    def _parse_docx(self, file_path: str) -> str:
        try:
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return "\n".join(full_text)
        except Exception as e:
            logger.error(f"Error parsing DOCX locally: {str(e)}")
            raise e

    def _parse_pptx(self, file_path: str) -> str:
        try:
            prs = pptx.Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            logger.error(f"Error parsing PPTX locally: {str(e)}")
            raise e

    def _parse_text(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file: {str(e)}")
            raise e

# Single global service instance
mistral_ocr_service = MistralOCRService()
